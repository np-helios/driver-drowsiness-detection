from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
SOURCE = ROOT / "Driver_Drowsiness_Research_Pitch_Premium.pptx"
OUTPUT = ROOT / "Driver_Drowsiness_Real_Visual_Pitch.pptx"

SCREEN_NORMAL = Path("/Users/nishthapandey/Desktop/Screenshot 2026-04-13 at 11.56.05 PM.png")
SCREEN_ALERT = Path("/Users/nishthapandey/Desktop/Screenshot 2026-04-13 at 11.56.17 PM.png")
SCREEN_TERMINAL = Path("/Users/nishthapandey/Desktop/Screenshot 2026-04-13 at 11.59.37 PM.png")
SCREEN_HOSTED = Path("/var/folders/jl/6fhfz73n5b73drgl2j4zsmnr0000gn/T/TemporaryItems/NSIRD_screencaptureui_W4vRcp/Screenshot 2026-04-13 at 11.56.50 PM.png")
SCREEN_CALIBRATION = Path("/var/folders/jl/6fhfz73n5b73drgl2j4zsmnr0000gn/T/TemporaryItems/NSIRD_screencaptureui_8eqq5X/Screenshot 2026-04-14 at 12.01.27 AM.png")


def main() -> None:
    prs = Presentation(str(SOURCE))

    # image-heavy slides: replace the main visual but keep the theme/caption areas
    replace_main_image(prs.slides[6], SCREEN_CALIBRATION)
    replace_main_image(prs.slides[8], SCREEN_ALERT)
    replace_main_image(prs.slides[10], SCREEN_HOSTED)
    replace_main_image(prs.slides[12], SCREEN_NORMAL)

    # bullet slides: add a screenshot to the right-side visual area without changing theme
    add_side_image(prs.slides[9], SCREEN_TERMINAL, "Execution snapshot")
    add_side_image(prs.slides[11], SCREEN_ALERT, "Alert-trigger example")

    prs.save(str(OUTPUT))
    print(f"Saved updated presentation to {OUTPUT}")


def replace_main_image(slide, image_path: Path) -> None:
    for shape in list(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            remove_shape(shape)
    add_framed_image(
        slide,
        image_path,
        left=Inches(1.0),
        top=Inches(1.47),
        width=Inches(9.15),
        height=Inches(5.0),
        frame_left=Inches(0.8),
        frame_top=Inches(1.25),
        frame_width=Inches(9.55),
        frame_height=Inches(5.45),
    )


def add_side_image(slide, image_path: Path, label: str) -> None:
    # cover the existing right-side callout with a screenshot card but keep the overall template
    cover = slide.shapes.add_shape(
        1,  # rectangle
        Inches(9.18),
        Inches(1.28),
        Inches(3.42),
        Inches(5.2),
    )
    cover.fill.solid()
    cover.fill.fore_color.rgb = RGBColor(14, 23, 38)
    cover.line.color.rgb = RGBColor(39, 58, 86)

    slide.shapes.add_picture(str(image_path), Inches(9.33), Inches(1.44), width=Inches(3.12), height=Inches(4.5))

    tag = slide.shapes.add_textbox(Inches(9.35), Inches(6.0), Inches(3.0), Inches(0.35))
    p = tag.text_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(251, 191, 36)


def add_framed_image(slide, image_path: Path, *, left, top, width, height, frame_left, frame_top, frame_width, frame_height) -> None:
    frame = slide.shapes.add_shape(
        1,  # rectangle
        frame_left,
        frame_top,
        frame_width,
        frame_height,
    )
    frame.fill.solid()
    frame.fill.fore_color.rgb = RGBColor(13, 19, 32)
    frame.line.color.rgb = RGBColor(39, 58, 86)
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def remove_shape(shape) -> None:
    sp = shape._element
    sp.getparent().remove(sp)


if __name__ == "__main__":
    main()
