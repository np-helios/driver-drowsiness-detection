from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
FIGS = ROOT / "report_figures"
OUT = ROOT / "Driver_Drowsiness_Research_Pitch_Premium.pptx"


BG = RGBColor(6, 10, 18)
SURFACE = RGBColor(14, 23, 38)
SURFACE_ALT = RGBColor(21, 32, 51)
NAVY = RGBColor(241, 245, 249)
SLATE = RGBColor(148, 163, 184)
BLUE = RGBColor(56, 189, 248)
GREEN = RGBColor(74, 222, 128)
AMBER = RGBColor(251, 191, 36)
RED = RGBColor(248, 113, 113)
PURPLE = RGBColor(167, 139, 250)
CARD = RGBColor(13, 19, 32)
CARD_EDGE = RGBColor(39, 58, 86)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Adaptive Driver Drowsiness Detection\nUsing Personalized Baseline Modeling",
        "A real-time, explainable, and deployment-aware driver monitoring framework",
    )
    add_bullets_slide(
        prs,
        "Problem Context",
        [
            "Driver drowsiness is a major safety risk that reduces reaction time and concentration.",
            "Vision-based monitoring is non-invasive, low-cost, and suitable for continuous observation.",
            "Most existing low-cost systems still rely on static thresholds shared by all users.",
        ],
        highlight="Key issue: generic thresholds do not reflect natural differences between drivers.",
    )
    add_bullets_slide(
        prs,
        "Research Gap",
        [
            "Natural blink rate, eye openness, and fatigue expression vary from person to person.",
            "Fixed EAR thresholds often create false positives and false negatives.",
            "Few practical systems combine personalization, explainability, and deployment awareness.",
        ],
        highlight="We need a personalized, real-time, and deployable drowsiness detection framework.",
    )
    add_bullets_slide(
        prs,
        "Proposed Solution",
        [
            "Learn the driver’s normal behavior during an initial calibration period.",
            "Track EAR, MAR, and 3D head pose continuously in real time.",
            "Compute fatigue through deviation from the driver’s own baseline instead of universal rules.",
            "Generate alerts, break recommendations, and runtime logs.",
        ],
        highlight="Core idea: drowsiness is personal, so detection should be personal too.",
    )
    add_image_slide(
        prs,
        "System Pipeline",
        "End-to-end architecture of the adaptive monitoring framework.",
        FIGS / "figure_01_system_architecture.png",
    )
    add_bullets_slide(
        prs,
        "Feature Extraction",
        [
            "EAR measures eye openness and sustained eye closure.",
            "MAR captures mouth opening and yawn-like behavior.",
            "3D head pose adds yaw, pitch, and roll as additional fatigue cues.",
            "Multi-signal monitoring is more robust than single-feature systems.",
        ],
        highlight="The detector combines eye, mouth, and head behavior into one interpretable pipeline.",
    )
    add_image_slide(
        prs,
        "Personalized Baseline Modeling",
        "The system first learns the driver and then monitors deviation from that learned profile.",
        FIGS / "figure_05_baseline_workflow.png",
    )
    add_image_slide(
        prs,
        "Deviation-Based Fatigue Scoring",
        "Fatigue is inferred from abnormal change relative to the individual’s baseline.",
        FIGS / "figure_06_deviation_scoring.png",
    )
    add_image_slide(
        prs,
        "Alert and Yawn Validation Logic",
        "Normal short blinks and casual mouth opening are filtered before triggering an alert.",
        FIGS / "figure_08_alert_flow.png",
    )
    add_bullets_slide(
        prs,
        "Implementation Highlights",
        [
            "Migrated the runtime from dlib to MediaPipe for modern compatibility.",
            "Refactored the monolithic prototype into a modular Python package.",
            "Added baseline persistence, event logging, session analytics, and deployment paths.",
            "Improved yawn logic to reduce false positives from speaking.",
        ],
        highlight="The project evolved from a script prototype into a structured, research-ready system.",
    )
    add_image_slide(
        prs,
        "Deployment Story",
        "Local runtime reflects practical edge deployment, while the hosted app demonstrates accessibility.",
        FIGS / "figure_07_deployment_comparison.png",
    )
    add_bullets_slide(
        prs,
        "Current Results",
        [
            "Real-time local monitoring runs successfully with live facial boxes and fatigue scoring.",
            "The system builds and reuses a per-driver baseline profile.",
            "Alerts trigger only for sustained abnormal behavior rather than normal blinks.",
            "Hosted deployment is available for public demonstration, though with browser-server latency.",
        ],
        highlight="The strongest real-time behavior is achieved in the local edge-style runtime.",
    )
    add_image_slide(
        prs,
        "Live Output Interface",
        "Representative interface of the final real-time local monitoring system.",
        FIGS / "figure_09_ui_mockup.png",
    )
    add_image_slide(
        prs,
        "Research Contribution",
        "Summary of the technical contributions implemented in the project.",
        FIGS / "figure_10_contribution_summary.png",
    )
    add_bullets_slide(
        prs,
        "Limitations and Future Scope",
        [
            "No large-scale multi-user evaluation has been completed yet.",
            "Hosted browser deployment is slower than native local execution.",
            "Future work includes gaze estimation, PERCLOS, low-light robustness, and embedded edge deployment.",
            "A comparative study against fixed-threshold systems would strengthen the research contribution further.",
        ],
        highlight="The current system is a strong prototype and a solid base for research-paper-level extension.",
    )
    add_bullets_slide(
        prs,
        "Conclusion",
        [
            "Driver drowsiness detection should not rely on universal blink rules.",
            "Personalized baseline modeling improves realism, explainability, and technical defensibility.",
            "This project demonstrates a practical pipeline combining EAR, MAR, head pose, and adaptive scoring.",
            "The local runtime is closest to real edge deployment, while the hosted version shows public deployability.",
        ],
        highlight="Final takeaway: adaptive personalization is a stronger foundation than one-size-fits-all thresholding.",
    )
    add_end_slide(prs, "Thank You", "Questions and Discussion")

    prs.save(OUT)
    print(f"Saved presentation to {OUT}")


def style_slide(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.18))
    band.fill.solid()
    band.fill.fore_color.rgb = BLUE
    band.line.fill.background()
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(10.95), Inches(0), Inches(2.383), Inches(7.5))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(10, 15, 27)
    panel.line.fill.background()
    glow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(6.95), Inches(13.333), Inches(0.55))
    glow.fill.solid()
    glow.fill.fore_color.rgb = RGBColor(8, 14, 24)
    glow.line.fill.background()


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style_slide(slide)
    accent_card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.9), Inches(5.8), Inches(0.55))
    accent_card.fill.solid()
    accent_card.fill.fore_color.rgb = SURFACE_ALT
    accent_card.line.color.rgb = CARD_EDGE
    at = accent_card.text_frame.paragraphs[0]
    at.text = "RESEARCH PITCH DECK"
    at.font.size = Pt(13)
    at.font.bold = True
    at.font.color.rgb = BLUE
    at.alignment = PP_ALIGN.CENTER
    tx = slide.shapes.add_textbox(Inches(0.8), Inches(1.65), Inches(8.6), Inches(2.3))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY
    sub = slide.shapes.add_textbox(Inches(0.85), Inches(4.15), Inches(7.4), Inches(1.0))
    sp = sub.text_frame.paragraphs[0]
    sp.text = subtitle
    sp.font.size = Pt(16)
    sp.font.color.rgb = SLATE
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.82), Inches(5.35), Inches(4.8), Inches(1.3))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD
    card.line.color.rgb = CARD_EDGE
    tf = card.text_frame
    tf.paragraphs[0].text = "Nishtha Pandey"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = NAVY
    p2 = tf.add_paragraph()
    p2.text = "B.Tech CSE | Manipal University Jaipur"
    p2.font.size = Pt(12)
    p2.font.color.rgb = SLATE
    quote = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.35), Inches(1.2), Inches(2.95), Inches(4.9))
    quote.fill.solid()
    quote.fill.fore_color.rgb = SURFACE
    quote.line.color.rgb = CARD_EDGE
    qtf = quote.text_frame
    qp = qtf.paragraphs[0]
    qp.text = "Most drowsiness systems assume a universal blink rule."
    qp.font.size = Pt(18)
    qp.font.bold = True
    qp.font.color.rgb = NAVY
    qp.alignment = PP_ALIGN.CENTER
    qp2 = qtf.add_paragraph()
    qp2.text = "This project replaces that assumption with personalized baseline modeling."
    qp2.font.size = Pt(15)
    qp2.font.color.rgb = BLUE
    qp2.alignment = PP_ALIGN.CENTER


def add_bullets_slide(prs: Presentation, title_text: str, bullets: list[str], highlight: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style_slide(slide)
    add_title(slide, title_text)
    body = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.35), Inches(8.45), Inches(5.15))
    body.fill.solid()
    body.fill.fore_color.rgb = CARD
    body.line.color.rgb = CARD_EDGE
    tf = body.text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(19)
        p.font.color.rgb = NAVY
        p.space_after = Pt(10)
        p.bullet = True
    side = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.35), Inches(1.35), Inches(2.6), Inches(1.05))
    side.fill.solid()
    side.fill.fore_color.rgb = SURFACE_ALT
    side.line.color.rgb = CARD_EDGE
    st = side.text_frame.paragraphs[0]
    st.text = "PITCH TAKEAWAY"
    st.font.size = Pt(13)
    st.font.bold = True
    st.font.color.rgb = AMBER
    st.alignment = PP_ALIGN.CENTER
    callout = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.25), Inches(2.65), Inches(3.25), Inches(3.85))
    callout.fill.solid()
    callout.fill.fore_color.rgb = SURFACE
    callout.line.color.rgb = CARD_EDGE
    ct = callout.text_frame
    ct.word_wrap = True
    p = ct.paragraphs[0]
    p.text = highlight
    p.font.size = Pt(19)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER
    stripe = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(9.25), Inches(2.65), Inches(0.1), Inches(3.85))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = BLUE
    stripe.line.fill.background()


def add_image_slide(prs: Presentation, title_text: str, caption: str, image_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style_slide(slide)
    add_title(slide, title_text)
    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.25), Inches(9.55), Inches(5.45))
    frame.fill.solid()
    frame.fill.fore_color.rgb = CARD
    frame.line.color.rgb = CARD_EDGE
    slide.shapes.add_picture(str(image_path), Inches(1.0), Inches(1.47), width=Inches(9.15), height=Inches(5.0))
    cap = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(10.55), Inches(1.45), Inches(2.0), Inches(4.15))
    cap.fill.solid()
    cap.fill.fore_color.rgb = SURFACE
    cap.line.color.rgb = CARD_EDGE
    tf = cap.text_frame
    p = tf.paragraphs[0]
    p.text = caption
    p.font.size = Pt(17)
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Research-pitch visual"
    p2.font.size = Pt(12)
    p2.font.color.rgb = BLUE
    p2.alignment = PP_ALIGN.CENTER
    badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(10.55), Inches(5.95), Inches(2.0), Inches(0.55))
    badge.fill.solid()
    badge.fill.fore_color.rgb = SURFACE_ALT
    badge.line.color.rgb = CARD_EDGE
    bp = badge.text_frame.paragraphs[0]
    bp.text = "KEY DIAGRAM"
    bp.font.size = Pt(12)
    bp.font.bold = True
    bp.font.color.rgb = AMBER
    bp.alignment = PP_ALIGN.CENTER


def add_end_slide(prs: Presentation, title_text: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style_slide(slide)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(7, 12, 22)
    band.line.fill.background()
    orb = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(9.0), Inches(0.7), Inches(3.0), Inches(3.0))
    orb.fill.solid()
    orb.fill.fore_color.rgb = RGBColor(16, 30, 54)
    orb.line.color.rgb = BLUE
    tx = slide.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(11.3), Inches(1.3))
    p = tx.text_frame.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    sub = slide.shapes.add_textbox(Inches(1.0), Inches(3.5), Inches(11.3), Inches(0.8))
    sp = sub.text_frame.paragraphs[0]
    sp.text = subtitle
    sp.font.size = Pt(18)
    sp.font.color.rgb = RGBColor(191, 219, 254)
    sp.alignment = PP_ALIGN.CENTER


def add_title(slide, text: str) -> None:
    marker = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.78), Inches(0.38), Inches(0.42), Inches(0.18))
    marker.fill.solid()
    marker.fill.fore_color.rgb = BLUE
    marker.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(0.78), Inches(0.62), Inches(10.1), Inches(0.7))
    p = tx.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = NAVY


if __name__ == "__main__":
    main()
